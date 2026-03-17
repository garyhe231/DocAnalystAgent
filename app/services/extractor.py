"""
Document text extractor — returns structured pages list.
Each page: {"page": int, "text": str, "sections": [{"heading": str, "lines": [{"line": int, "text": str}]}]}
"""
import re
from typing import List, Dict, Any, Optional
from pathlib import Path


def extract(file_path: str, filename: str) -> Dict[str, Any]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(file_path)
    elif suffix in (".docx",):
        return _extract_docx(file_path)
    elif suffix in (".pptx",):
        return _extract_pptx(file_path)
    elif suffix in (".xlsx", ".xls"):
        return _extract_xlsx(file_path)
    elif suffix in (".txt", ".md", ".csv"):
        return _extract_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


def _build_sections(raw_text: str) -> List[Dict[str, Any]]:
    """Split raw text into line-numbered sections."""
    lines = raw_text.split("\n")
    sections: List[Dict[str, Any]] = []
    current_heading = "Content"
    current_lines: List[Dict[str, Any]] = []
    global_line = 1

    heading_re = re.compile(r"^(#{1,6}\s+.+|[A-Z][A-Z0-9 ,\-:]{3,60}$|\d+[\.\)]\s+.+)")

    for raw_line in lines:
        stripped = raw_line.strip()
        if not stripped:
            global_line += 1
            continue
        is_heading = bool(heading_re.match(stripped)) and len(stripped) < 120
        if is_heading and current_lines:
            sections.append({"heading": current_heading, "lines": current_lines})
            current_heading = stripped
            current_lines = []
        else:
            current_lines.append({"line": global_line, "text": stripped})
        global_line += 1

    if current_lines:
        sections.append({"heading": current_heading, "lines": current_lines})
    return sections


def _extract_pdf(file_path: str) -> Dict[str, Any]:
    import fitz  # PyMuPDF
    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text")
        pages.append({
            "page": i + 1,
            "text": text,
            "sections": _build_sections(text),
        })
    doc.close()
    return {"pages": pages, "total_pages": len(pages)}


def _extract_docx(file_path: str) -> Dict[str, Any]:
    from docx import Document
    doc = Document(file_path)
    # Group paragraphs into virtual pages of ~50 paragraphs
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    page_size = 50
    pages = []
    for i in range(0, max(len(paras), 1), page_size):
        chunk = paras[i:i + page_size]
        text = "\n".join(chunk)
        pages.append({
            "page": len(pages) + 1,
            "text": text,
            "sections": _build_sections(text),
        })
    if not pages:
        pages = [{"page": 1, "text": "", "sections": []}]
    return {"pages": pages, "total_pages": len(pages)}


def _extract_pptx(file_path: str) -> Dict[str, Any]:
    from pptx import Presentation
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        text = "\n".join(parts)
        pages.append({
            "page": i + 1,
            "text": text,
            "sections": _build_sections(text),
        })
    if not pages:
        pages = [{"page": 1, "text": "", "sections": []}]
    return {"pages": pages, "total_pages": len(pages)}


def _extract_xlsx(file_path: str) -> Dict[str, Any]:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    pages = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(c) if c is not None else "" for c in row)
            if row_text.strip():
                rows.append(row_text)
        text = f"Sheet: {sheet.title}\n" + "\n".join(rows)
        pages.append({
            "page": len(pages) + 1,
            "text": text,
            "sections": _build_sections(text),
        })
    if not pages:
        pages = [{"page": 1, "text": "", "sections": []}]
    return {"pages": pages, "total_pages": len(pages)}


def _extract_txt(file_path: str) -> Dict[str, Any]:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    page_size = 3000  # chars per virtual page
    pages = []
    for i in range(0, max(len(content), 1), page_size):
        text = content[i:i + page_size]
        pages.append({
            "page": len(pages) + 1,
            "text": text,
            "sections": _build_sections(text),
        })
    if not pages:
        pages = [{"page": 1, "text": "", "sections": []}]
    return {"pages": pages, "total_pages": len(pages)}


def full_text(doc_data: Dict[str, Any]) -> str:
    """Flatten all pages into a single string with page markers."""
    parts = []
    for p in doc_data["pages"]:
        parts.append(f"\n=== Page {p['page']} ===\n{p['text']}")
    return "\n".join(parts)
